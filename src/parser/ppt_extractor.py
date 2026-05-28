"""PPT/PPTX 文本提取器 — 无需 LibreOffice，纯 Python 实现。

- .pptx: python-pptx（文本、表格、备注）
- .ppt: olefile + 二进制记录解析（提取 Unicode 文本和 ASCII 文本）
"""

import logging
import struct
from pathlib import Path

from .base import ParserExecutionError

logger = logging.getLogger(__name__)


# PowerPoint 模板占位文本，需过滤
_TEMPLATE_TEXTS = [
    "单击此处编辑母版标题样式",
    "单击此处编辑母版文本样式",
    "单击此处添加标题",
    "单击此处添加文本",
    "Click to edit Master title style",
    "Click to edit Master text styles",
    "Click to edit Master subtitle style",
    "Second level",
    "Third level",
    "Fourth level",
    "Fifth level",
    "第二级",
    "第三级",
    "第四级",
    "第五级",
]


def _is_template_line(line: str) -> bool:
    """检查一行是否为 PPT 模板占位文本。"""
    stripped = line.strip()
    if not stripped:
        return True
    for tmpl in _TEMPLATE_TEXTS:
        if stripped == tmpl or stripped.startswith(tmpl):
            return True
    return False


def _filter_template(text: str) -> str:
    """移除模板占位行和空幻灯片。"""
    blocks = text.split("\n\n")
    filtered: list[str] = []
    for block in blocks:
        lines = block.strip().split("\n")
        # 保留非模板行
        kept = [line for line in lines if not _is_template_line(line)]
        # 移除只有幻灯片标题标记的空块（如只有 "## 幻灯片 N"）
        real_content = [line for line in kept if not line.startswith("## 幻灯片")]
        if real_content:
            filtered.append("\n".join(kept))
    return "\n\n".join(filtered)


def extract_text(file_path: Path) -> str:
    """统一入口：根据扩展名选择提取方法。"""
    ext = file_path.suffix.lower()
    if ext == ".pptx":
        text = _extract_pptx_text(file_path)
    elif ext == ".ppt":
        text = _extract_ppt_text(file_path)
    else:
        raise ParserExecutionError(f"不支持的文件格式: {ext}")

    return _filter_template(text)


def _extract_pptx_text(file_path: Path) -> str:
    """从 .pptx 文件提取所有文本（python-pptx）。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ParserExecutionError("python-pptx 未安装。请运行: pip install python-pptx")

    prs = Presentation(str(file_path))
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        slide_lines.append(f"\n## 幻灯片 {slide_idx}\n")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_lines.append(row_text)

        # 备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"\n> 备注: {notes}")

        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))

    return "\n\n".join(parts)


# ---- 二进制 .ppt 解析 ----

# MS-PPT record types that contain text
_RT_TEXT_CHARS_ATOM = 0x0FA0  # Unicode text (UTF-16LE)
_RT_TEXT_BYTES_ATOM = 0x0FA8  # ASCII/ANSI text

# Record header: 2 bytes version+instance, 2 bytes recType, 4 bytes length
_RECORD_HEADER = struct.Struct("<HH I")


def _extract_ppt_text(file_path: Path) -> str:
    """从旧格式 .ppt 文件提取文本（OLE2 二进制扫描）。

    PPT 二进制格式使用嵌套记录结构，文本存储在 TextCharsAtom (0x0FA0)
    记录中。此函数扫描整个流查找文本记录，无需完整解析记录层级。
    """
    try:
        import olefile
    except ImportError:
        raise ParserExecutionError("olefile 未安装。请运行: pip install olefile")

    if not olefile.isOleFile(str(file_path)):
        raise ParserExecutionError(f"不是有效的 OLE2/PPT 文件: {file_path.name}")

    ole = olefile.OleFileIO(str(file_path))
    try:
        ppt_stream = ole.openstream("PowerPoint Document")
        raw_bytes = ppt_stream.read()
    finally:
        ole.close()

    text_parts: list[str] = []
    i = 0
    n = len(raw_bytes)

    while i < n - 8:
        # 搜索 TextCharsAtom 标记 (0x0FA0 = A0 0F little-endian)
        marker_pos = raw_bytes.find(b"\xa0\x0f", i)
        if marker_pos == -1:
            break

        # 验证这是一个有效的记录头
        if marker_pos + 8 > n:
            break

        _ver_inst = int.from_bytes(raw_bytes[marker_pos - 2 : marker_pos], "little")  # noqa: F841
        rec_len = int.from_bytes(raw_bytes[marker_pos + 2 : marker_pos + 6], "little")

        # 合理性检查：文本记录通常 < 10000 字节，> 0 字节
        if 0 < rec_len < 10000 and marker_pos + 6 + rec_len <= n:
            rec_data = raw_bytes[marker_pos + 6 : marker_pos + 6 + rec_len]
            try:
                text = rec_data.decode("utf-16-le")
                # 验证文本质量：需要包含可打印字符
                if text and any(c.isalnum() or "一" <= c <= "鿿" for c in text):
                    text = text.replace("\r\n", "\n").replace("\r", "\n")
                    text = text.strip("\x00\n ")
                    if text and len(text) > 1:
                        text_parts.append(text)
            except UnicodeDecodeError:
                pass

        i = marker_pos + 6

    # Also scan for TextBytesAtom (0x0FA8) for ANSI text
    i = 0
    while i < n - 8:
        marker_pos = raw_bytes.find(b"\xa8\x0f", i)
        if marker_pos == -1:
            break
        if marker_pos + 8 > n:
            break

        rec_len = int.from_bytes(raw_bytes[marker_pos + 2 : marker_pos + 6], "little")
        if 0 < rec_len < 10000 and marker_pos + 6 + rec_len <= n:
            rec_data = raw_bytes[marker_pos + 6 : marker_pos + 6 + rec_len]
            for encoding in ["gbk", "utf-8"]:
                try:
                    text = rec_data.decode(encoding)
                    if any(c.isalnum() or "一" <= c <= "鿿" for c in text):
                        text = text.replace("\r\n", "\n").replace("\r", "\n")
                        text = text.strip("\x00\n ")
                        if text and len(text) > 1:
                            text_parts.append(text)
                    break
                except UnicodeDecodeError:
                    continue
        i = marker_pos + 6

    if not text_parts:
        return ""

    # 去重（相邻相同文本只保留一次）
    deduped: list[str] = []
    for text in text_parts:
        if not deduped or text != deduped[-1]:
            deduped.append(text)
    text_parts = deduped

    # 分组为幻灯片（启发式：短标题文本标记新幻灯片开始）
    result_lines: list[str] = []
    slide_count = 0
    current_slide: list[str] = []

    for text in text_parts:
        if not current_slide:
            slide_count += 1
            current_slide.append(f"\n## 幻灯片 {slide_count}\n")
            current_slide.append(text)
        elif len(text) < 80 and any(
            kw in text for kw in ["第", "章", "Chapter", "节", "课", "讲", "§", "信号", "系统"]
        ):
            result_lines.append("\n".join(current_slide))
            slide_count += 1
            current_slide = [f"\n## 幻灯片 {slide_count}\n", text]
        else:
            current_slide.append(text)

    if current_slide:
        result_lines.append("\n".join(current_slide))

    return "\n\n".join(result_lines)
